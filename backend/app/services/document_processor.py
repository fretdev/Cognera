"""
Document text extraction — PDF, Word, PowerPoint, plain text, CSV, Markdown.
"""
import io
import fitz  # PyMuPDF


def extract_text(file_bytes: bytes, content_type: str, filename: str = "") -> str:
    ct   = (content_type or "").lower()
    name = (filename or "").lower()

    if ct == "application/pdf" or name.endswith(".pdf"):
        return _extract_pdf(file_bytes)

    if ct in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ) or name.endswith((".docx", ".doc")):
        return _extract_docx(file_bytes)

    if ct in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ) or name.endswith((".pptx", ".ppt")):
        return _extract_pptx(file_bytes)

    if ct.startswith("text/") or name.endswith((".txt", ".md", ".csv", ".markdown")):
        return file_bytes.decode("utf-8", errors="replace")

    try:
        return file_bytes.decode("utf-8", errors="replace")
    except Exception:
        raise ValueError(f"Unsupported file type: {content_type}")


def _extract_pdf(pdf_bytes: bytes) -> str:
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        return "\n\n".join(page.get_text() for page in doc)


def _extract_docx(docx_bytes: bytes) -> str:
    from docx import Document
    doc   = Document(io.BytesIO(docx_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _extract_pptx(pptx_bytes: bytes) -> str:
    from pptx import Presentation
    prs   = Presentation(io.BytesIO(pptx_bytes))
    parts = []
    for num, slide in enumerate(prs.slides, 1):
        texts = [f"[Slide {num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if len(texts) > 1:
            parts.append("\n".join(texts))
    return "\n\n".join(parts)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    words = text.split()
    if not words:
        return []
    chunks, start = [], 0
    while start < len(words):
        chunk = " ".join(words[start : start + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks
